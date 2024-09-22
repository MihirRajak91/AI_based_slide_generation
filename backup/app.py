import streamlit as st
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import extract_sections
import mistral_summarizer
import fitz  
import re

MAX_BULLETS_PER_SLIDE = 6  # Maximum number of bullet points per slide

# Function to create PowerPoint presentation with summarized text and images
def create_presentation(filename):
    # Extract sections, images, and text
    sections_and_images = extract_sections.extract_sections_and_images(filename)
    
    summarized_sections = {}
    
    # Loop through each section and send the text to Mistral for summarization
    for section, content in sections_and_images.items():
        summary = mistral_summarizer.mistral_summarize(content['text'])
        
        if summary:
            summarized_sections[section] = {
                "text": summary, 
                "images": content.get("images", [])
            }
        else:
            summarized_sections[section] = {
                "text": "No summary available",
                "images": content.get("images", [])
            }

    prs = Presentation()

    #  function to add a slide with bullet points
    def add_slide_with_text(prs, title, bullet_points, is_first_slide=True):
        slide_layout = prs.slide_layouts[1]  # Layout for title and content
        slide = prs.slides.add_slide(slide_layout)

        if is_first_slide:
            slide.shapes.title.text = title
            # Formatting the title
            title_shape = slide.shapes.title
            title_text_frame = title_shape.text_frame
            title_paragraph = title_text_frame.paragraphs[0]
            title_paragraph.font.size = Pt(30)  # Set title size
            title_paragraph.font.bold = True  # Bold title
            title_paragraph.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue color for the title

        text_frame = slide.placeholders[1].text_frame
        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(20)

    # Helper function to split bullet points if they exceed max per slide
    def split_bullet_points(bullet_points, max_points=MAX_BULLETS_PER_SLIDE):
        return [bullet_points[i:i + max_points] for i in range(0, len(bullet_points), max_points)]

    # Helper function to add image slide
    def add_image_slide(prs, img_path):
        slide_layout = prs.slide_layouts[6]  # Blank layout for full image
        slide = prs.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(5.5)
        slide.shapes.add_picture(img_path, left, top, width, height)

    # Loop over summarized sections and create slides
    for section_title, content in summarized_sections.items():
        # Split the summarized text into bullet points
        bullet_points = [re.sub(r'(Topic:|Summary:)', '', point).strip() for point in content["text"].split('\n') if point.strip()]
        bullet_point_chunks = split_bullet_points(bullet_points)

        # Add slides with bullet points (multiple if necessary)
        for idx, bullet_chunk in enumerate(bullet_point_chunks):
            add_slide_with_text(prs, section_title, bullet_chunk, is_first_slide=(idx == 0))

        # Add images, if any, for this section
        if content["images"]:
            for img_path in content["images"]:
                add_image_slide(prs, img_path)

    # Save the presentation to a file-like object
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    return ppt_io

# Function to display the PDF file
def display_pdf(file):
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        image = page.get_pixmap()  # Render the page as an image
        img_data = image.pil_tobytes("jpeg")  # Convert to bytes to render in Streamlit
        st.image(img_data, caption=f"Page {page_num + 1}", use_column_width=True)

# Streamlit app front end
st.title("PDF to PowerPoint Generator")

uploaded_pdf = st.file_uploader("Upload a PDF", type="pdf")

if uploaded_pdf is not None:
    st.subheader("Uploaded PDF Preview")
    
    # Preview the PDF in the app
    with st.spinner("Loading PDF..."):
        display_pdf(uploaded_pdf)
    
    # Process the PDF and generate the presentation
    st.subheader("Generating Presentation...")
    with st.spinner("Processing..."):
        ppt_io = create_presentation(uploaded_pdf)

    # Provide download option for the generated PowerPoint
    st.success("Presentation Generated!")
    st.download_button(
        label="Download Presentation",
        data=ppt_io,
        file_name="presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
