from pptx import Presentation
from pptx.util import Inches, Pt
import re

MAX_BULLETS_PER_SLIDE = 6  # Maximum number of bullet points per slide

def create_slide_without_images(prs, title, bullet_points, selected_font):
    """
    Creates a slide without images, using the placeholders for title and content.
    :param prs: Presentation object.
    :param title: Slide title.
    :param bullet_points: List of bullet points.
    :param selected_font: Font to be used for the text.
    """
# Use a blank slide layout to have more control over positioning
    slide_layout = prs.slide_layouts[5]  # Using blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the title
    title_placeholder = slide.shapes.title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Create a new text box for the bullet points and place it on the slide
    left = Inches(0.75)  # Adjust left position for text box
    top = Inches(2.34)   # Adjust top position
    width = Inches(12)  # Adjust width for text box
    height = Inches(5)   # Adjust height for text box

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Add bullet points inside the text box
    for i, point in enumerate(bullet_points):
        if i == 0:
            # First paragraph uses the default paragraph
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(19)  # Set font size
        p.font.name = selected_font
        p.level = 0  # Set bullet point level


def create_slide_with_single_image(prs, title, bullet_points, selected_font, img_path):
    """
    Creates a slide with a title, bullet points, and a single image with customized positions.
    :param prs: PowerPoint presentation object.
    :param title: Slide title.
    :param bullet_points: List of bullet points.
    :param selected_font: Font to be used for the text.
    :param img_path: Path to the image file.
    """
    # Use a blank slide layout to have more control over positioning
    slide_layout = prs.slide_layouts[5]  # Using blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the title
    title_placeholder = slide.shapes.title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Create a new text box for the bullet points and place it on the slide
    left = Inches(0.75)  # Adjust left position for text box
    top = Inches(2.34)   # Adjust top position
    width = Inches(6.66)  # Adjust width for text box
    height = Inches(3.99)   # Adjust height for text box

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Add bullet points inside the text box
    for i, point in enumerate(bullet_points):
        if i == 0:
            # First paragraph uses the default paragraph
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(19)  # Set font size
        p.font.name = selected_font
        p.level = 0  # Set bullet point level

    # Add the image on the right side of the slide
    img_left = Inches(8.61)  # Adjusted position to match layout
    img_top = Inches(3.13)
    img_width = Inches(3)
    img_height = Inches(3)
    slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)


def create_slide_with_two_images(prs, title, bullet_points, selected_font, img_path1, img_path2):
    """
    Creates a slide with a title, bullet points, and two images using placeholders.
    :param prs: PowerPoint presentation object.
    :param title: Slide title.
    :param bullet_points: List of bullet points.
    :param selected_font: Font to be used for the text.
    :param img_path1: Path to the first image file.
    :param img_path2: Path to the second image file.
    """
    slide_layout = prs.slide_layouts[1]  # Layout with title and content placeholders
    slide = prs.slides.add_slide(slide_layout)

    # Set the title using the title placeholder
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Set bullet points using the content placeholder
    content_placeholder = slide.placeholders[1].text_frame
    for point in bullet_points:
        p = content_placeholder.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.name = selected_font

    # Add the first image
    left = Inches(6.5)
    top = Inches(1.5)
    width = Inches(3)
    height = Inches(3)
    slide.shapes.add_picture(img_path1, left, top, width, height)

    # Add the second image
    left = Inches(6.5)
    top = Inches(4.5)
    width = Inches(3)
    height = Inches(3)
    slide.shapes.add_picture(img_path2, left, top, width, height)

def add_image_slide(prs, img_path):
    """
    Add a slide with a single image.
    :param prs: PowerPoint presentation object.
    :param img_path: Path to the image file.
    """
    # Add a blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout for blank slide

    # Define the image's position and size
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5.5)

    # Add the image to the slide
    slide.shapes.add_picture(img_path, left, top, width, height)

def create_presentation(prs, summarized_dict, images_dict, selected_font):
    """
    Creates a PowerPoint presentation from the summarized dictionary and adds images.
    Handles different cases: no images, single image, two images, and more.
    :param prs: PowerPoint presentation object from template.
    :param summarized_dict: Dictionary with section titles as keys and bullet points as values.
    :param images_dict: Dictionary with section titles as keys and image paths as values.
    :param selected_font: Selected font for text.
    """
    def split_bullet_points(bullet_points, max_points=MAX_BULLETS_PER_SLIDE):
        """
        Splits the bullet points into chunks of max_points per chunk.
        """
        return [bullet_points[i:i + max_points] for i in range(0, len(bullet_points), max_points)]

    # Iterate over summarized sections
    for section_title, bullet_points in summarized_dict.items():
        cleaned_bullet_points = [re.sub(r'(Topic:|Summary:)', '', point).strip() for point in bullet_points.split('\n') if point.strip()]
        bullet_point_chunks = split_bullet_points(cleaned_bullet_points)

        # Check if the section exists in images_dict and how many images it has
        images = images_dict.get(section_title, None)

        if images is None or len(images) == 0:  # No images
            for idx, bullet_chunk in enumerate(bullet_point_chunks):
                create_slide_without_images(prs, section_title, bullet_chunk, selected_font)
        elif len(images) == 1:  # One image
            for idx, bullet_chunk in enumerate(bullet_point_chunks):
                create_slide_with_single_image(prs, section_title, bullet_chunk, selected_font, images[0])
        elif len(images) == 2:  # Two images
            for idx, bullet_chunk in enumerate(bullet_point_chunks):
                create_slide_with_two_images(prs, section_title, bullet_chunk, selected_font, images[0], images[1])
        else:  # More than two images
            for img_path in images:
                add_image_slide(prs, img_path)

    # # Save the presentation
    # prs.save("presentation_output.pptx")
    # print("Presentation saved as 'presentation_output.pptx'")

# Example of calling the create_presentation function:
if __name__ == "__main__":
    prs = Presentation()

    # Sample data for demonstration
    summarized_dict = {
        "Introduction": "This is the introduction.\nHere are the main points.",
        "Conclusion": "This is the conclusion.\nThese are the key takeaways."
    }

    images_dict = {
        "Introduction": ["path/to/image1.png"],
        "Conclusion": ["path/to/image2.png", "path/to/image3.png"]
    }

    selected_font = "Arial"

    # Create the presentation
    create_presentation(prs, summarized_dict, images_dict, selected_font)
