import streamlit as st
import os
import pdfplumber
from pptx import Presentation
from extract_sections import extract_sections_and_images
from summarize_sections import summarize_sections, send_to_mistral_for_bullet_points
from pptx_exp import create_presentation
from PIL import Image

# Define available fonts and templates
font_choices = ["Arial", "Calibri", "Times New Roman", "Verdana", "Georgia"]

# Path where templates are stored
TEMPLATE_DIR = "D:/python/Projects/slide_generation/last_hope/templates"  # Change this to your actual template directory
UPLOAD_DIR = "uploads"  # Directory for saving uploaded PDFs

# List of available PPTX templates
pptx_files = [f for f in os.listdir(TEMPLATE_DIR) if f.endswith('.pptx')]

# Function to convert PDF pages to images and display them
def display_pdf_as_images(pdf_path):
    images = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            st.write(f"Displaying page {i + 1}...")
            page_image = page.to_image(resolution=300)
            img_path = os.path.join(UPLOAD_DIR, f"page_{i+1}.png")
            page_image.save(img_path)
            images.append(img_path)
            st.image(img_path, caption=f"Page {i + 1}", use_column_width=True)
    return images

# Function to update the title in the first slide
def update_presentation_title(prs, new_title):
    # Access the first slide (usually the title slide)
    first_slide = prs.slides[0]
    # Access the title placeholder (placeholder 0)
    title_placeholder = first_slide.shapes.title
    if title_placeholder is not None:
        title_placeholder.text = new_title

# Streamlit frontend
def main():
    st.title("PDF to PowerPoint Generator")

    # Step 1: Input for presentation name (to replace title on first slide)
    presentation_name = st.text_input("Enter the name of your presentation:")

    # Step 2: Upload PDF file
    uploaded_pdf = st.file_uploader("Upload your PDF file", type="pdf")

    # Step 3: Choose a PPTX template
    template_choice = st.selectbox("Choose a PPTX template", pptx_files)

    # Step 4: Choose a font
    font_choice = st.selectbox("Choose a font", font_choices)

    # Display the Start button
    generate_presentation = st.button("Start Presentation Generation")

    if uploaded_pdf is not None:
        # Ensure the 'uploads' directory exists
        if not os.path.exists(UPLOAD_DIR):
            os.makedirs(UPLOAD_DIR)

        # Save the uploaded PDF
        pdf_path = os.path.join(UPLOAD_DIR, uploaded_pdf.name)
        with open(pdf_path, "wb") as f:
            f.write(uploaded_pdf.getbuffer())

        # Convert PDF pages to images and display them
        st.write("Displaying the uploaded PDF as images...")
        pdf_images = display_pdf_as_images(pdf_path)

        # Generate the presentation only when the button is pressed
        if generate_presentation:
            st.write("Starting presentation generation process...")

            # Step 1: Extract sections and images from the PDF
            st.write("Extracting sections and images from the PDF...")
            content_dict = extract_sections_and_images(pdf_path)
            st.write("Extraction complete.")
            st.write(content_dict)  # Display the extracted content

            # Step 2: Summarize the sections
            st.write("Summarizing sections...")
            summarized_dict = summarize_sections(content_dict)
            st.write("Summarization complete.")

            # Step 3: Generate bullet points for sections
            st.write("Generating bullet points...")
            bullet_point_dict = send_to_mistral_for_bullet_points(summarized_dict)
            st.write("Bullet point generation complete.")

            # Step 4: Prepare images dictionary
            st.write("Preparing image data...")
            images_dict = {section: content.get("images", []) for section, content in content_dict.items()}

            # Step 5: Load the chosen PowerPoint template
            ppt_template_path = os.path.join(TEMPLATE_DIR, template_choice)
            prs = Presentation(ppt_template_path)

            # Step 6: Update the title on the first slide with the presentation name
            if presentation_name:
                st.write(f"Updating title on the first slide to: {presentation_name}")
                update_presentation_title(prs, presentation_name)

            # Step 7: Generate the presentation
            st.write("Generating the PowerPoint presentation...")
            create_presentation(prs, bullet_point_dict, images_dict, font_choice)
            st.write("Presentation generation complete.")

            # Save the generated presentation
            output_ppt_path = f"{presentation_name}.pptx"
            prs.save(output_ppt_path)

            # Button to download the generated presentation
            st.write("Your presentation is ready for download.")
            with open(output_ppt_path, "rb") as f:
                st.download_button(
                    label="Download Presentation",
                    data=f,
                    file_name=output_ppt_path,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

if __name__ == "__main__":
    main()
